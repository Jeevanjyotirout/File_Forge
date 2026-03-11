import os
import uuid
import io
import zipfile
from typing import List

from PIL import Image
import img2pdf
from pypdf import PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt

from app.config import settings


def _out_path(filename: str) -> str:
    return os.path.join(settings.OUTPUT_DIR, filename)


def _unique(prefix: str, ext: str) -> str:
    return f"{prefix}_{uuid.uuid4().hex[:8]}.{ext}"


def _extract_pdf_text_by_page(pdf_path: str) -> List[str]:
    reader = PdfReader(pdf_path)
    return [page.extract_text() or "" for page in reader.pages]


# ─── PDF → WORD ──────────────────────────────────────────────────────────────

def pdf_to_word(file_path: str) -> str:
    """
    Convert PDF to .docx by extracting text and creating a Word document.
    For pixel-perfect conversion, pdf2docx is used when available.
    """
    out_name = _unique("converted", "docx")
    out_path = _out_path(out_name)

    try:
        from pdf2docx import Converter
        cv = Converter(file_path)
        cv.convert(out_path)
        cv.close()
    except Exception:
        # Fallback: plain text extraction into docx
        from docx import Document
        pages = _extract_pdf_text_by_page(file_path)
        doc = Document()
        doc.add_heading("Converted Document", 0)
        for i, text in enumerate(pages):
            doc.add_heading(f"Page {i + 1}", level=2)
            doc.add_paragraph(text)
        doc.save(out_path)

    return out_path


# ─── PDF → EXCEL ─────────────────────────────────────────────────────────────

def pdf_to_excel(file_path: str) -> str:
    """Extract PDF text into an Excel workbook, one sheet per page."""
    pages = _extract_pdf_text_by_page(file_path)
    out_name = _unique("converted", "xlsx")
    out_path = _out_path(out_name)

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # remove default sheet

    for i, text in enumerate(pages):
        ws = wb.create_sheet(title=f"Page {i + 1}")
        lines = text.splitlines()
        for row_idx, line in enumerate(lines, start=1):
            # Try to split by common delimiters for table-like data
            parts = line.split("\t") if "\t" in line else [line]
            for col_idx, part in enumerate(parts, start=1):
                ws.cell(row=row_idx, column=col_idx, value=part.strip())

    wb.save(out_path)
    return out_path


# ─── PDF → PPT ───────────────────────────────────────────────────────────────

def pdf_to_ppt(file_path: str) -> str:
    """Convert each PDF page into a PowerPoint slide with extracted text."""
    pages = _extract_pdf_text_by_page(file_path)
    out_name = _unique("converted", "pptx")
    out_path = _out_path(out_name)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    for i, text in enumerate(pages):
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Title line
        p = tf.paragraphs[0]
        p.text = f"Page {i + 1}"
        p.font.bold = True
        p.font.size = Pt(18)

        # Content
        if text.strip():
            from pptx.util import Pt
            from pptx.oxml.ns import qn
            p2 = tf.add_paragraph()
            p2.text = text[:1000]  # cap per slide
            p2.font.size = Pt(11)

    prs.save(out_path)
    return out_path


# ─── JPG/IMAGE → PDF ─────────────────────────────────────────────────────────

def images_to_pdf(file_paths: List[str]) -> str:
    """Combine one or more images (JPG, PNG, BMP, etc.) into a single PDF."""
    out_name = _unique("images", "pdf")
    out_path = _out_path(out_name)

    # Convert all images to RGB JPEG bytes for img2pdf
    processed: List[bytes] = []
    for fp in file_paths:
        img = Image.open(fp).convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=90)
        processed.append(buf.getvalue())

    with open(out_path, "wb") as f:
        f.write(img2pdf.convert(processed))

    return out_path


# ─── BATCH HELPER ────────────────────────────────────────────────────────────

def batch_convert(file_paths: List[str], target_format: str) -> str:
    """Convert multiple files and return a ZIP of results."""
    converters = {
        "docx": pdf_to_word,
        "xlsx": pdf_to_excel,
        "pptx": pdf_to_ppt,
    }
    fn = converters.get(target_format)
    if not fn:
        raise ValueError(f"Unsupported target format: {target_format}")

    out_files = [fn(fp) for fp in file_paths]

    zip_name = _unique("batch_converted", "zip")
    zip_path = _out_path(zip_name)
    with zipfile.ZipFile(zip_path, "w") as zf:
        for f in out_files:
            zf.write(f, arcname=os.path.basename(f))

    return zip_path
